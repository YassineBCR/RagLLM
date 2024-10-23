from langchain_community.document_loaders import PyPDFLoader, TextLoader
import os
from typing import List
from langchain_core.documents import Document
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.vectorstores import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter
import subprocess
import json
import whisper
from docx import Document as DocxDocument  # Pour charger les fichiers Word
from openpyxl import load_workbook  # Pour charger les fichiers Excel
from pptx import Presentation  # Pour charger les fichiers PowerPoint

def get_ollama_models() -> List[str]:
    """
    Récupère la liste des modèles disponibles sur Ollama.

    Returns:
        List[str]: Liste des noms de modèles disponibles.
    """
    try:
        result = subprocess.run(['ollama', 'list', '--format=json'], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        if result.returncode != 0:
            print(f"Erreur lors de l'exécution de 'ollama list': {result.stderr}")
            return []
        models_output = result.stdout
        models_data = json.loads(models_output)
        models = [model['name'] for model in models_data]
        return models
    except Exception as e:
        print(f"Exception lors de la récupération des modèles Ollama: {e}")
        return []

def load_document_into_database(model_name: str, document_path: str, chunk_size: int = 1000, chunk_overlap: int = 100) -> Chroma:
    """
    Loads a single document into the Chroma database after splitting the text into chunks.

    Args:
        model_name (str): The name of the embedding model to use.
        document_path (str): The path to the document file to load.
        chunk_size (int): The size of the chunks to split the document into.
        chunk_overlap (int): The overlap between chunks.

    Returns:
        Chroma: The Chroma database with loaded documents.
    """
    print(f"Loading document: {document_path}")
    raw_document = load_document(document_path)  # Appelle la fonction load_document
    
    # Utiliser les valeurs de chunk_size et chunk_overlap pour créer un TextSplitter
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    documents = text_splitter.split_documents([raw_document])

    print("Creating embeddings and loading document into Chroma")
    embeddings = OllamaEmbeddings(model=model_name)
    db = Chroma.from_documents(
        documents,
        embeddings,
    )
    return db

# Fonction pour transcrire un fichier audio en texte
def transcribe_audio_to_text(audio_file_path, output_txt_path):
    model = whisper.load_model("base")  # Utilisez le modèle "base" de Whisper pour la transcription
    result = model.transcribe(audio_file_path)
    
    # Sauvegarder le texte transcrit dans un fichier .txt
    with open(output_txt_path, 'w', encoding='utf-8') as f:
        f.write(result['text'])
    
    return output_txt_path

def load_document(file_path: str) -> Document:
    """
    Loads a single document from the specified file path based on its extension.

    Args:
        file_path (str): The path to the file to load.

    Returns:
        Document: A loaded document object with necessary metadata.

    Raises:
        FileNotFoundError: If the specified file does not exist.
        ValueError: If the file extension is not supported.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"The specified file does not exist: {file_path}")

    # Vérifiez si c'est un fichier audio et effectuez la transcription
    if file_path.endswith(".mp3") or file_path.endswith(".wav"):
        print(f"Transcribing audio file: {file_path}")
        txt_file_path = os.path.splitext(file_path)[0] + ".txt"  # Créez un fichier .txt pour la transcription
        transcribe_audio_to_text(file_path, txt_file_path)
        file_path = txt_file_path  # Remplacez le chemin du fichier par celui du fichier transcrit

    # Initialiser `document` avec une valeur par défaut
    document = []

    # Charger en fonction de l'extension de fichier
    if file_path.endswith(".pdf"):
        print(f"Loading PDF file: {file_path}")
        loader = PyPDFLoader(file_path)
        document = loader.load()
    elif file_path.endswith(".md") or file_path.endswith(".txt"):
        print(f"Loading text/markdown file: {file_path}")
        loader = TextLoader(file_path)
        document = loader.load()
    elif file_path.endswith(".docx"):
        print(f"Loading Word file: {file_path}")
        document = load_word_file(file_path)
    elif file_path.endswith(".xlsx"):
        print(f"Loading Excel file: {file_path}")
        document = load_excel_file(file_path)
    elif file_path.endswith(".pptx"):
        print(f"Loading PowerPoint file: {file_path}")
        document = load_powerpoint_file(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_path}")

    # Ajoutez les métadonnées manquantes pour les fichiers
    for doc in document:
        if "page" not in doc.metadata:
            doc.metadata["page"] = 1  # Valeur par défaut
        if "source" not in doc.metadata:
            doc.metadata["source"] = os.path.basename(file_path)  # Nom du fichier comme source

    # Vérifiez que `document` est bien initialisé avant de le retourner
    if not document:
        raise ValueError(f"Erreur lors du chargement du document depuis '{file_path}'.")

    return document[0]  

def load_word_file(file_path: str) -> List[Document]:
    """
    Load a Word (.docx) file and convert it into a Document object.

    Args:
        file_path (str): Path to the Word file.

    Returns:
        List[Document]: A list containing the Document object with extracted text.
    """
    doc = DocxDocument(file_path)
    full_text = '\n'.join([para.text for para in doc.paragraphs])
    return [Document(page_content=full_text, metadata={"source": os.path.basename(file_path)})]

def load_excel_file(file_path: str) -> List[Document]:
    """
    Load an Excel (.xlsx) file and convert its content into a Document object.

    Args:
        file_path (str): Path to the Excel file.

    Returns:
        List[Document]: A list containing the Document object with extracted text.
    """
    wb = load_workbook(file_path)
    full_text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        sheet_data = []
        for row in ws.iter_rows(values_only=True):
            sheet_data.append('\t'.join([str(cell) for cell in row if cell is not None]))
        full_text.append(f"Sheet: {sheet}\n" + '\n'.join(sheet_data))
    
    return [Document(page_content='\n\n'.join(full_text), metadata={"source": os.path.basename(file_path)})]

def load_powerpoint_file(file_path: str) -> List[Document]:
    """
    Load a PowerPoint (.pptx) file and convert its content into a Document object.

    Args:
        file_path (str): Path to the PowerPoint file.

    Returns:
        List[Document]: A list containing the Document object with extracted text.
    """
    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        full_text.append('\n'.join(slide_text))
    
    return [Document(page_content='\n\n'.join(full_text), metadata={"source": os.path.basename(file_path)})]
